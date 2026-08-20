"""
Generate the hackathon pitch-deck presentation.
Run: python scripts/generate_ppt.py
Output: docs/Predictive_Maintenance_Agent_Pitch.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pathlib import Path
import math

# ── Colour palette (dark industrial-tech) ──────────────────────────
BG_DARK   = RGBColor(0x12, 0x18, 0x26)   # deep navy-charcoal
BG_CARD   = RGBColor(0x1C, 0x24, 0x38)   # slightly lighter card bg
TEAL      = RGBColor(0x00, 0xD4, 0xAA)   # electric teal accent
AMBER     = RGBColor(0xFF, 0xA0, 0x2F)   # amber for highlights
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xCC, 0xD0, 0xD8)   # muted body text
DIM       = RGBColor(0x7A, 0x84, 0x9E)   # dimmed / secondary
RED       = RGBColor(0xFF, 0x4D, 0x4D)
GREEN     = RGBColor(0x00, 0xE6, 0x76)
DARK_TEXT  = RGBColor(0x0A, 0x0E, 0x1A)

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

# ── Helpers ────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]   # completely blank layout

def _add_bg(slide, color=BG_DARK):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color

def _rect(slide, left, top, w, h, fill=None, border=None, opacity=None):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, w, h)
    shape.line.fill.background()
    if border:
        shape.line.color.rgb = border
        shape.line.width = Pt(1.5)
    if fill:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill
    else:
        shape.fill.background()
    shape.shadow.inherit = False
    return shape

def _circle(slide, left, top, size, fill=TEAL):
    shape = slide.shapes.add_shape(MSO_SHAPE.OVAL, left, top, size, size)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def _text_box(slide, left, top, w, h, text, size=18, color=WHITE,
              bold=False, align=PP_ALIGN.LEFT, font_name="Calibri"):
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return txBox

def _multi_text(slide, left, top, w, h, lines, default_size=18,
                default_color=LIGHT, line_spacing=1.5):
    """lines = list of (text, size, color, bold, align) or just str."""
    txBox = slide.shapes.add_textbox(left, top, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(lines):
        if isinstance(item, str):
            txt, sz, clr, bld, al = item, default_size, default_color, False, PP_ALIGN.LEFT
        else:
            txt = item[0]
            sz  = item[1] if len(item) > 1 else default_size
            clr = item[2] if len(item) > 2 else default_color
            bld = item[3] if len(item) > 3 else False
            al  = item[4] if len(item) > 4 else PP_ALIGN.LEFT
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = txt
        p.font.size = Pt(sz)
        p.font.color.rgb = clr
        p.font.bold  = bld
        p.font.name  = "Calibri"
        p.alignment  = al
        p.space_after = Pt(4)
    return txBox

def _accent_line(slide, left, top, width, color=TEAL):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, Pt(4))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.shadow.inherit = False
    return shape

def _footer(slide, text="Cognizant Hackathon 2026  |  Predictive Maintenance Agent"):
    _text_box(slide, Inches(0.5), Inches(7.0), Inches(12), Inches(0.4),
              text, size=10, color=DIM, align=PP_ALIGN.CENTER)

def _slide_number(slide, num):
    _text_box(slide, Inches(12.3), Inches(7.0), Inches(0.8), Inches(0.4),
              str(num), size=10, color=DIM, align=PP_ALIGN.RIGHT)

def _card(slide, left, top, w, h, title, body_lines, icon_text="",
          fill=BG_CARD, border=TEAL, title_color=TEAL):
    _rect(slide, left, top, w, h, fill=fill, border=border)
    if icon_text:
        _text_box(slide, left + Inches(0.25), top + Inches(0.15),
                  Inches(0.5), Inches(0.5), icon_text, size=24, color=TEAL)
    t_left = left + Inches(0.25) + (Inches(0.6) if icon_text else Inches(0))
    _text_box(slide, t_left, top + Inches(0.15),
              w - Inches(0.6), Inches(0.4), title, size=16, color=title_color, bold=True)
    _multi_text(slide, left + Inches(0.3), top + Inches(0.6),
                w - Inches(0.6), h - Inches(0.8), body_lines,
                default_size=13, default_color=LIGHT)

def _metric_box(slide, left, top, w, h, value, label, val_color=TEAL):
    _rect(slide, left, top, w, h, fill=BG_CARD, border=TEAL)
    _text_box(slide, left, top + Inches(0.2), w, Inches(0.6),
              value, size=36, color=val_color, bold=True, align=PP_ALIGN.CENTER)
    _text_box(slide, left, top + Inches(0.85), w, Inches(0.4),
              label, size=12, color=LIGHT, align=PP_ALIGN.CENTER)

# ────────────────────────────────────────────────────────────────────
# SLIDE 1 — Title
# ────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
_add_bg(s)

# decorative circles (abstract factory gauge)
for cx, cy, sz, clr in [
    (Inches(10.5), Inches(0.8), Inches(1.8), RGBColor(0x00,0xD4,0xAA)),
    (Inches(11.0), Inches(2.2), Inches(1.0), RGBColor(0xFF,0xA0,0x2F)),
    (Inches(0.3),  Inches(5.5), Inches(1.2), RGBColor(0x00,0xD4,0xAA)),
]:
    c = _circle(s, cx, cy, sz, fill=clr)
    c.fill.fore_color.rgb = clr
    # make semi-transparent via alpha (approximate with lighter mix)

_accent_line(s, Inches(1.2), Inches(2.2), Inches(2.5), TEAL)

_text_box(s, Inches(1.2), Inches(2.5), Inches(10), Inches(1.5),
          "Predictive Maintenance\nAgent", size=48, color=WHITE, bold=True)

_text_box(s, Inches(1.2), Inches(4.2), Inches(9), Inches(0.8),
          "An AI that doesn't just predict machine failure —\nit explains why, and decides what to do about it.",
          size=22, color=LIGHT)

_text_box(s, Inches(1.2), Inches(5.5), Inches(6), Inches(0.5),
          "Cognizant Hackathon 2026", size=16, color=TEAL, bold=True)

_footer(s)

# ────────────────────────────────────────────────────────────────────
# SLIDE 2 — The Story (hook)
# ────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
_add_bg(s)

_text_box(s, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
          "The $19,600 Mistake That Never Had to Happen", size=36, color=WHITE, bold=True)
_accent_line(s, Inches(0.8), Inches(1.15), Inches(3), AMBER)

# Story flow — 4 cards
story_cards = [
    ("01", "Running Normally", "A milling machine runs\nfor 14 hours straight.\nNobody watches every second."),
    ("02", "Something Shifts", "Tool wear climbs.\nTorque creeps up.\nNobody notices."),
    ("03", "Breakdown", "The machine grinds to\na halt. Production stops.\n3 hours to diagnose."),
    ("04", "The Fix", "A worn-out tool bit.\n10-minute swap —\nif only someone had known."),
]
for i, (num, title, body) in enumerate(story_cards):
    x = Inches(0.8 + i * 3.1)
    _rect(s, x, Inches(1.8), Inches(2.8), Inches(3.2), fill=BG_CARD, border=DIM)
    _text_box(s, x + Inches(0.2), Inches(1.95), Inches(0.6), Inches(0.5),
              num, size=28, color=TEAL, bold=True)
    _text_box(s, x + Inches(0.2), Inches(2.5), Inches(2.4), Inches(0.4),
              title, size=18, color=WHITE, bold=True)
    _text_box(s, x + Inches(0.2), Inches(3.0), Inches(2.4), Inches(1.8),
              body, size=14, color=LIGHT)

# Pivot line
_rect(s, Inches(0.8), Inches(5.5), Inches(11.7), Inches(1.2), fill=RGBColor(0x0A,0x2A,0x3A), border=TEAL)
_text_box(s, Inches(1.2), Inches(5.7), Inches(11), Inches(0.8),
          "What if the machine could tell you — before it breaks — exactly what's wrong, why, and what to do?",
          size=22, color=TEAL, bold=True, align=PP_ALIGN.CENTER)

_footer(s, "Cognizant Hackathon 2026  |  Predictive Maintenance Agent")

# ────────────────────────────────────────────────────────────────────
# SLIDE 3 — The Core Idea
# ────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
_add_bg(s)

_text_box(s, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
          "A Doctor for Machines", size=36, color=WHITE, bold=True)
_accent_line(s, Inches(0.8), Inches(1.15), Inches(2.5), TEAL)

_text_box(s, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
          "Think of it as a doctor that never sleeps — and always shows its reasoning.",
          size=18, color=DIM)

capabilities = [
    ("Watches", "Continuously monitors vital signs:\ntemperature, speed, torque, tool wear", "01"),
    ("Predicts", "Estimates failure risk —\nlike a health risk score at the doctor", "02"),
    ("Diagnoses", "Identifies the exact cause —\nnot 'something's wrong' but 'this specific part'", "03"),
    ("Decides", "Keep running, watch closely,\nor stop the line now", "04"),
]
for i, (title, body, num) in enumerate(capabilities):
    x = Inches(0.8 + i * 3.1)
    _card(s, x, Inches(2.3), Inches(2.8), Inches(3.0),
          f"{num}. {title}", [body], fill=BG_CARD, border=TEAL)

_footer(s)

# ────────────────────────────────────────────────────────────────────
# SLIDE 4 — Why This Is Hard
# ────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
_add_bg(s)

_text_box(s, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
          "The Problem Nobody Talks About", size=36, color=WHITE, bold=True)
_accent_line(s, Inches(0.8), Inches(1.15), Inches(3), RED)

# Left: accuracy paradox
_card(s, Inches(0.8), Inches(1.8), Inches(5.5), Inches(2.5),
      "The Accuracy Trap", [
          "Machines fail only ~3.4% of the time",
          "A lazy AI that always says 'fine' scores 96.6%",
          "That's a perfect accuracy score — and it catches ZERO failures",
          ("Accuracy is meaningless on imbalanced data.", 13, AMBER, True),
      ], icon_text="01", border=RED)

# Right: black box
_card(s, Inches(6.8), Inches(1.8), Inches(5.7), Inches(2.5),
      "The Black-Box Problem", [
          "ML models say '87% risk' — then nothing else",
          "A technician needs to know WHAT to fix",
          "'High risk' without a cause = useless to the floor",
          ("Explainability isn't a bonus — it's the entire point.", 13, AMBER, True),
      ], icon_text="02", border=AMBER)

# Bottom callout
_rect(s, Inches(0.8), Inches(4.8), Inches(11.7), Inches(1.5), fill=RGBColor(0x0A,0x2A,0x3A), border=TEAL)
_multi_text(s, Inches(1.3), Inches(5.0), Inches(11), Inches(1.2), [
    ("Our Bet:", 20, TEAL, True),
    ("We built a system where every prediction comes with a physics-backed explanation.", 18, WHITE),
    ("The agent can't just say 'risky' — it has to prove WHY, against documented engineering limits.", 16, LIGHT),
])

_footer(s)

# ────────────────────────────────────────────────────────────────────
# SLIDE 5 — Agent vs Dashboard
# ────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
_add_bg(s)

_text_box(s, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
          "It Decides, Not Just Displays", size=36, color=WHITE, bold=True)
_accent_line(s, Inches(0.8), Inches(1.15), Inches(2.5), TEAL)

# Left: Dashboard (faded)
_rect(s, Inches(0.8), Inches(1.8), Inches(5.5), Inches(4.0), fill=BG_CARD, border=DIM)
_text_box(s, Inches(1.2), Inches(1.95), Inches(4.5), Inches(0.4),
          "Traditional Dashboard", size=20, color=DIM, bold=True)
_multi_text(s, Inches(1.2), Inches(2.5), Inches(4.8), Inches(3.0), [
    ("Shows risk score", 16, DIM),
    ("'Something's wrong'", 16, DIM),
    ("Human interprets the number", 16, DIM),
    ("No action taken by the system", 16, DIM),
    ("", 10, DIM),
    ("Passive  —  shows data, decides nothing", 14, DIM, True),
])

# Right: Agent (bright)
_rect(s, Inches(6.8), Inches(1.8), Inches(5.7), Inches(4.0), fill=BG_CARD, border=TEAL)
_text_box(s, Inches(7.2), Inches(1.95), Inches(4.5), Inches(0.4),
          "Our Agent", size=20, color=TEAL, bold=True)
_multi_text(s, Inches(7.2), Inches(2.5), Inches(5.0), Inches(3.0), [
    ("Makes a decision", 16, WHITE),
    ("'Overstrain failure — stop now'", 16, WHITE),
    ("Agent recommends action", 16, WHITE),
    ("Structured work order with evidence", 16, WHITE),
    ("", 10, WHITE),
    ("Active  —  reasons, decides, acts", 14, TEAL, True),
])

# Bottom: 3 decisions
_text_box(s, Inches(0.8), Inches(6.1), Inches(11.7), Inches(0.4),
          "Three Real Decisions:", size=18, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

decisions = [
    ("ESCALATE NOW", "Stop the line — verified critical failure", RED),
    ("SCHEDULE", "Fix at next planned stop", AMBER),
    ("MONITOR", "Continue watching — nominal reading", GREEN),
]
for i, (label, desc, clr) in enumerate(decisions):
    x = Inches(0.8 + i * 4.1)
    _rect(s, x, Inches(6.5), Inches(3.8), Inches(0.5), fill=None, border=clr)
    _text_box(s, x + Inches(0.15), Inches(6.52), Inches(3.5), Inches(0.45),
              f"{label}  —  {desc}", size=12, color=clr, bold=True, align=PP_ALIGN.CENTER)

_footer(s)

# ────────────────────────────────────────────────────────────────────
# SLIDE 6 — Architecture
# ────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
_add_bg(s)

_text_box(s, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
          "How It Thinks", size=36, color=WHITE, bold=True)
_accent_line(s, Inches(0.8), Inches(1.15), Inches(2), TEAL)

steps = [
    ("1", "SENSE", "Reads live sensor data:\ntemperature, speed,\ntorque, tool wear", TEAL),
    ("2", "PREDICT", "ML model estimates\nfailure probability\nusing XGBoost", AMBER),
    ("3", "VERIFY", "Physics rules engine\nchecks against known\nengineering limits", GREEN),
    ("4", "DECIDE", "Combines both opinions;\nonly acts when\nthey agree", RGBColor(0xFF, 0x66, 0x99)),
]

for i, (num, title, body, clr) in enumerate(steps):
    x = Inches(0.8 + i * 3.15)
    _rect(s, x, Inches(1.8), Inches(2.8), Inches(3.2), fill=BG_CARD, border=clr)
    _circle(s, x + Inches(1.0), Inches(2.0), Inches(0.8), fill=clr)
    _text_box(s, x + Inches(1.0), Inches(2.05), Inches(0.8), Inches(0.7),
              num, size=28, color=BG_DARK, bold=True, align=PP_ALIGN.CENTER)
    _text_box(s, x + Inches(0.2), Inches(3.0), Inches(2.4), Inches(0.4),
              title, size=18, color=clr, bold=True, align=PP_ALIGN.CENTER)
    _text_box(s, x + Inches(0.2), Inches(3.5), Inches(2.4), Inches(1.2),
              body, size=13, color=LIGHT, align=PP_ALIGN.CENTER)

# Arrows between steps
for i in range(3):
    x = Inches(3.55 + i * 3.15)
    _text_box(s, x, Inches(2.15), Inches(0.5), Inches(0.5),
              ">", size=28, color=DIM, bold=True, align=PP_ALIGN.CENTER)

# Callout
_rect(s, Inches(0.8), Inches(5.3), Inches(11.7), Inches(1.2), fill=RGBColor(0x0A,0x2A,0x3A), border=TEAL)
_text_box(s, Inches(1.3), Inches(5.5), Inches(11), Inches(0.8),
          "If the AI and the physics disagree, the agent flags the conflict — it never guesses a cause it can't prove.",
          size=18, color=TEAL, bold=True, align=PP_ALIGN.CENTER)

_footer(s)

# ────────────────────────────────────────────────────────────────────
# SLIDE 7 — Tech Stack
# ────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
_add_bg(s)

_text_box(s, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
          "Built to Be Defensible", size=36, color=WHITE, bold=True)
_accent_line(s, Inches(0.8), Inches(1.15), Inches(2), TEAL)

stack_groups = [
    ("Data & Modeling", [
        "Python, XGBoost, IsolationForest",
        "Trained on AI4I 2020 dataset",
        "10,000 machine-cycle readings",
        "80/20 stratified split, tuned threshold",
    ], TEAL),
    ("Explainability", [
        "SHAP (TreeExplainer)",
        "Feature-by-feature attribution",
        "Global + per-prediction breakdown",
        "Top features match physics rules",
    ], AMBER),
    ("Reasoning Layer", [
        "LLM agent with tool-calling",
        "7 dedicated tools for evidence",
        "Deterministic offline fallback",
        "Structured work orders (Pydantic)",
    ], GREEN),
    ("Knowledge + Interface", [
        "Maintenance playbook (JSON)",
        "Cost model with cited rates",
        "Live dashboard (Streamlit)",
        "Chat interface for 'why?' questions",
    ], RGBColor(0xFF, 0x66, 0x99)),
]

for i, (group_title, items, clr) in enumerate(stack_groups):
    x = Inches(0.8 + i * 3.15)
    _rect(s, x, Inches(1.8), Inches(2.8), Inches(4.2), fill=BG_CARD, border=clr)
    _text_box(s, x + Inches(0.2), Inches(1.95), Inches(2.4), Inches(0.4),
              group_title, size=16, color=clr, bold=True)
    _accent_line(s, x + Inches(0.2), Inches(2.35), Inches(2.0), clr)
    for j, item in enumerate(items):
        _text_box(s, x + Inches(0.3), Inches(2.6 + j * 0.55), Inches(2.3), Inches(0.5),
                  item, size=13, color=LIGHT)

_footer(s)

# ────────────────────────────────────────────────────────────────────
# SLIDE 8 — Demo Narrative
# ────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
_add_bg(s)

_text_box(s, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
          "See It In Action", size=36, color=WHITE, bold=True)
_accent_line(s, Inches(0.8), Inches(1.15), Inches(2), TEAL)

beats = [
    ("1", "Normal", "Machine running.\nDashboard green.\nAll sensors nominal.", GREEN),
    ("2", "Warning", "Tool wear climbing.\nTorque rising.\nRisk score increasing.", AMBER),
    ("3", "Alert", "Physics rule fires.\nSHAP shows top driver.\nRoot cause identified.", RED),
    ("4", "Decision", "Work order issued.\nSeverity: CRITICAL.\nAction: Stop the line.", TEAL),
]

for i, (num, title, body, clr) in enumerate(beats):
    x = Inches(0.8 + i * 3.15)
    _rect(s, x, Inches(1.8), Inches(2.8), Inches(3.5), fill=BG_CARD, border=clr)
    _circle(s, x + Inches(0.15), Inches(2.0), Inches(0.5), fill=clr)
    _text_box(s, x + Inches(0.15), Inches(2.02), Inches(0.5), Inches(0.45),
              num, size=18, color=BG_DARK, bold=True, align=PP_ALIGN.CENTER)
    _text_box(s, x + Inches(0.75), Inches(2.0), Inches(1.8), Inches(0.4),
              title, size=18, color=WHITE, bold=True)
    _text_box(s, x + Inches(0.25), Inches(2.7), Inches(2.3), Inches(2.2),
              body, size=14, color=LIGHT)

# Demo case callouts
_rect(s, Inches(0.8), Inches(5.7), Inches(11.7), Inches(1.2), fill=BG_CARD, border=TEAL)
_multi_text(s, Inches(1.3), Inches(5.8), Inches(11), Inches(1.0), [
    ("Demo Case — UDI 70 (Overstrain Failure)", 16, TEAL, True),
    ("wear x torque = 12,549 minNm vs 11,000 limit (114%). SHAP top feature: osf_margin +4.86.", 14, LIGHT),
    ("Agent decision: ESCALATE NOW. Cost avoided: ~$19,600. Confidence: HIGH (ML + physics agree).", 14, WHITE),
])

_footer(s)

# ────────────────────────────────────────────────────────────────────
# SLIDE 9 — Results
# ────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
_add_bg(s)

_text_box(s, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
          "Metrics That Actually Matter", size=36, color=WHITE, bold=True)
_accent_line(s, Inches(0.8), Inches(1.15), Inches(2.5), AMBER)

# Big metric boxes
metrics = [
    ("0.88", "PR-AUC\n(26x baseline)", TEAL),
    ("87%", "Recall\n59/68 failures caught", GREEN),
    ("79%", "Precision\n16 false alarms", AMBER),
    ("100%", "Root-Cause Accuracy\nat HIGH confidence", RGBColor(0xFF, 0x66, 0x99)),
]

for i, (val, label, clr) in enumerate(metrics):
    x = Inches(0.8 + i * 3.15)
    _metric_box(s, x, Inches(1.8), Inches(2.8), Inches(1.8), val, label, val_color=clr)

# Table
_rect(s, Inches(0.8), Inches(4.0), Inches(11.7), Inches(2.2), fill=BG_CARD, border=DIM)
headers = ["Metric", "Value", "What It Means"]
rows = [
    ["PR-AUC", "0.88 (vs 0.034 baseline)", "The honest metric for rare-failure data"],
    ["Recall", "87% (59/68)", "We'd rather double-check a healthy machine than miss a real failure"],
    ["Precision", "79% (16 false alarms)", "Low false-alarm rate across 2,000 test cycles"],
    ["Root-Cause Accuracy", "100% at HIGH confidence", "When it names a cause, it's never wrong — structural guarantee"],
]
# Header row
for j, h in enumerate(headers):
    x = Inches(1.0 + j * 3.8)
    _text_box(s, x, Inches(4.15), Inches(3.5), Inches(0.35),
              h, size=13, color=TEAL, bold=True)
# Data rows
for i, row in enumerate(rows):
    for j, cell in enumerate(row):
        x = Inches(1.0 + j * 3.8)
        _text_box(s, x, Inches(4.55 + i * 0.42), Inches(3.5), Inches(0.35),
                  cell, size=12, color=LIGHT)

# Callout
_rect(s, Inches(0.8), Inches(6.5), Inches(11.7), Inches(0.5), fill=RGBColor(0x2A,0x15,0x15), border=RED)
_text_box(s, Inches(1.3), Inches(6.52), Inches(11), Inches(0.4),
          "We do NOT lead with accuracy — a model that always says 'fine' scores 96.6% and is useless.",
          size=13, color=RED, bold=True, align=PP_ALIGN.CENTER)

_footer(s)

# ────────────────────────────────────────────────────────────────────
# SLIDE 10 — Business Impact
# ────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
_add_bg(s)

_text_box(s, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
          "The Factory Bottom Line", size=36, color=WHITE, bold=True)
_accent_line(s, Inches(0.8), Inches(1.15), Inches(2), AMBER)

# Left: cost breakdown
_card(s, Inches(0.8), Inches(1.8), Inches(5.5), Inches(2.5),
      "Cost of Unplanned Downtime", [
          "$2,600/hr — mid-market CNC cell (cited rate)",
          "Emergency repairs cost 3-5x scheduled ones",
          "Cascading delays across the production line",
          "Labor idle time + lost margin",
      ], icon_text="$$", border=AMBER, title_color=AMBER)

# Right: our impact
_card(s, Inches(6.8), Inches(1.8), Inches(5.7), Inches(2.5),
      "What Our Agent Prevents", [
          "59 failures caught early across 2,000 cycles",
          "Emergency repairs become scheduled fixes",
          "2.5 hr avg downtime per failure avoided",
          ("$383,000 estimated savings across test set", 14, GREEN, True),
      ], icon_text="->", border=GREEN, title_color=GREEN)

# Before / After comparison
_text_box(s, Inches(0.8), Inches(4.6), Inches(11.7), Inches(0.5),
          "Before vs After", size=20, color=WHITE, bold=True, align=PP_ALIGN.CENTER)

_rect(s, Inches(0.8), Inches(5.1), Inches(5.5), Inches(1.6), fill=BG_CARD, border=RED)
_text_box(s, Inches(1.2), Inches(5.2), Inches(4.5), Inches(0.4),
          "Reactive Maintenance", size=16, color=RED, bold=True)
_multi_text(s, Inches(1.2), Inches(5.65), Inches(4.8), Inches(1.0), [
    ("Machine breaks without warning", 13, LIGHT),
    ("Emergency callout, overtime costs", 13, LIGHT),
    ("Collateral damage to other parts", 13, LIGHT),
    ("Unpredictable, expensive, unsafe", 13, RED, True),
])

_rect(s, Inches(6.8), Inches(5.1), Inches(5.7), Inches(1.6), fill=BG_CARD, border=GREEN)
_text_box(s, Inches(7.2), Inches(5.2), Inches(4.5), Inches(0.4),
          "With Our Agent", size=16, color=GREEN, bold=True)
_multi_text(s, Inches(7.2), Inches(5.65), Inches(5.0), Inches(1.0), [
    ("Early warning with verified cause", 13, LIGHT),
    ("Scheduled fix at next planned stop", 13, LIGHT),
    ("Minimal disruption to production", 13, LIGHT),
    ("Predictable, cheaper, safer", 13, GREEN, True),
])

_footer(s)

# ────────────────────────────────────────────────────────────────────
# SLIDE 11 — Roadmap
# ────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
_add_bg(s)

_text_box(s, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
          "From Prototype to Product", size=36, color=WHITE, bold=True)
_accent_line(s, Inches(0.8), Inches(1.15), Inches(2.5), TEAL)

phases = [
    ("NOW", "Hackathon Prototype", [
        "Single-machine prototype",
        "Physics-grounded synthetic data",
        "Offline + LLM agent paths",
        "Live interactive dashboard",
    ], TEAL),
    ("NEXT", "Pilot Deployment", [
        "Connect to real IoT sensors",
        "Multiple machines on one floor",
        "Integration with existing SCADA",
        "Mobile alerts for technicians",
    ], AMBER),
    ("LATER", "Fleet Scale", [
        "Multi-factory dashboard",
        "ERP maintenance integration",
        "Fleet-wide anomaly patterns",
        "Predictive parts inventory",
    ], GREEN),
    ("VISION", "Standard AI Layer", [
        "Plug-in for any manufacturer",
        "Works with any sensor suite",
        "Licensable platform",
        "Industry-standard explainable AI",
    ], RGBColor(0xFF, 0x66, 0x99)),
]

for i, (phase, title, items, clr) in enumerate(phases):
    x = Inches(0.8 + i * 3.15)
    _rect(s, x, Inches(1.8), Inches(2.8), Inches(4.5), fill=BG_CARD, border=clr)
    _text_box(s, x + Inches(0.2), Inches(1.95), Inches(2.4), Inches(0.4),
              phase, size=14, color=clr, bold=True)
    _text_box(s, x + Inches(0.2), Inches(2.35), Inches(2.4), Inches(0.4),
              title, size=16, color=WHITE, bold=True)
    _accent_line(s, x + Inches(0.2), Inches(2.75), Inches(1.8), clr)
    for j, item in enumerate(items):
        _text_box(s, x + Inches(0.3), Inches(3.0 + j * 0.55), Inches(2.3), Inches(0.5),
                  item, size=13, color=LIGHT)

_footer(s)

# ────────────────────────────────────────────────────────────────────
# SLIDE 12 — Limitations (honest)
# ────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
_add_bg(s)

_text_box(s, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
          "What We're Honest About", size=36, color=WHITE, bold=True)
_accent_line(s, Inches(0.8), Inches(1.15), Inches(2.5), AMBER)

_text_box(s, Inches(0.8), Inches(1.5), Inches(11), Inches(0.5),
          "Judges respect maturity. Here's what we know, and what we're doing about it.",
          size=16, color=DIM)

limitations = [
    ("Synthetic Data", "Built on the AI4I 2020 dataset — physically grounded, but not live plant data yet. Our next step is real IoT sensor feeds.", AMBER),
    ("No Real Timestamps", "Data is snapshot-based, not a true time series. UDI is used as a simulated replay index for the demo.", DIM),
    ("RNF is Irreducible", "One failure mode (Random) fires at 0.1% regardless of sensor values. We exclude it by design and say so.", RED),
    ("Probabilities Uncalibrated", "scale_pos_weight distorts raw probabilities. We fitted a calibrator — marginal improvement, recommend leaving it out.", DIM),
]

for i, (title, body, clr) in enumerate(limitations):
    y = Inches(2.2 + i * 1.2)
    _rect(s, Inches(0.8), y, Inches(11.7), Inches(1.0), fill=BG_CARD, border=clr)
    _text_box(s, Inches(1.2), y + Inches(0.1), Inches(2.5), Inches(0.35),
              title, size=16, color=clr, bold=True)
    _text_box(s, Inches(1.2), y + Inches(0.5), Inches(10.8), Inches(0.45),
              body, size=13, color=LIGHT)

_footer(s)

# ────────────────────────────────────────────────────────────────────
# SLIDE 13 — Why This Wins
# ────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
_add_bg(s)

_text_box(s, Inches(0.8), Inches(0.5), Inches(10), Inches(0.7),
          "Why This Wins", size=36, color=WHITE, bold=True)
_accent_line(s, Inches(0.8), Inches(1.15), Inches(2), TEAL)

reasons = [
    ("1", "Physics-Verified", "Every explanation is backed by\nreal engineering limits — not guesses", TEAL),
    ("2", "Honest Metrics", "PR-AUC over accuracy.\nWe report what matters.", AMBER),
    ("3", "Real Decisions", "Not a dashboard — an agent\nthat commits to an action", GREEN),
    ("4", "Clear Path to Product", "Prototype -> pilot -> fleet.\nThis is a licensable platform.", RGBColor(0xFF, 0x66, 0x99)),
    ("5", "Self-Aware", "We know our limitations\nand we state them upfront", WHITE),
]

for i, (num, title, body, clr) in enumerate(reasons):
    x = Inches(0.8 + i * 2.5)
    _rect(s, x, Inches(1.8), Inches(2.2), Inches(3.5), fill=BG_CARD, border=clr)
    _circle(s, x + Inches(0.7), Inches(2.0), Inches(0.7), fill=clr)
    _text_box(s, x + Inches(0.7), Inches(2.05), Inches(0.7), Inches(0.6),
              num, size=24, color=BG_DARK, bold=True, align=PP_ALIGN.CENTER)
    _text_box(s, x + Inches(0.2), Inches(2.9), Inches(1.8), Inches(0.4),
              title, size=15, color=clr, bold=True, align=PP_ALIGN.CENTER)
    _text_box(s, x + Inches(0.15), Inches(3.35), Inches(1.9), Inches(1.5),
              body, size=12, color=LIGHT, align=PP_ALIGN.CENTER)

_footer(s)

# ────────────────────────────────────────────────────────────────────
# SLIDE 14 — Thank You
# ────────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
_add_bg(s)

# Decorative elements
for cx, cy, sz, clr in [
    (Inches(9.5), Inches(0.5), Inches(2.0), TEAL),
    (Inches(10.5), Inches(2.0), Inches(1.2), AMBER),
    (Inches(0.5), Inches(5.0), Inches(1.5), GREEN),
]:
    _circle(s, cx, cy, sz, fill=clr)

_accent_line(s, Inches(1.2), Inches(2.2), Inches(2.5), TEAL)

_text_box(s, Inches(1.2), Inches(2.5), Inches(10), Inches(1.5),
          "Thank You", size=52, color=WHITE, bold=True)

_text_box(s, Inches(1.2), Inches(4.2), Inches(9), Inches(0.8),
          "\"We didn't just build a model that predicts failure —\nwe built a system that can defend its own reasoning.\"",
          size=22, color=TEAL, bold=True)

_text_box(s, Inches(1.2), Inches(5.5), Inches(6), Inches(0.5),
          "Cognizant Hackathon 2026", size=16, color=DIM)

_text_box(s, Inches(1.2), Inches(6.0), Inches(6), Inches(0.5),
          "Questions?", size=28, color=WHITE, bold=True)

_footer(s)

# ── Save ───────────────────────────────────────────────────────────
out = Path(__file__).resolve().parent.parent / "docs" / "Predictive_Maintenance_Agent_Pitch.pptx"
out.parent.mkdir(parents=True, exist_ok=True)
prs.save(str(out))
print(f"Saved: {out}")
print(f"Slides: {len(prs.slides)}")
